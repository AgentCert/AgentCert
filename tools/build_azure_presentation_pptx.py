from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("/mnt/d/Studies/AgentCert/AgentCert_Azure_Deployment_Presentation.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(50, 84, 168)
LIGHT_BLUE = RGBColor(237, 244, 255)
SKY = RGBColor(122, 168, 255)
GREEN = RGBColor(34, 139, 94)
LIGHT_GREEN = RGBColor(238, 251, 244)
ORANGE = RGBColor(230, 126, 34)
LIGHT_ORANGE = RGBColor(255, 247, 238)
PURPLE = RGBColor(110, 72, 170)
LIGHT_PURPLE = RGBColor(246, 240, 255)
GRAY = RGBColor(93, 106, 125)
LIGHT_GRAY = RGBColor(245, 247, 250)
DARK = RGBColor(33, 43, 54)
WHITE = RGBColor(255, 255, 255)


def add_title(slide, title: str, subtitle: str | None = None, number: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(9.8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.9), Inches(9.8), Inches(0.35))
        sub = sub_box.text_frame.paragraphs[0]
        sub.text = subtitle
        sub.font.size = Pt(11)
        sub.font.color.rgb = GRAY

    if number:
        num_box = slide.shapes.add_textbox(Inches(12.15), Inches(0.22), Inches(0.8), Inches(0.3))
        p = num_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = number
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY


def add_footer(slide, text: str = "AgentCert Azure Production Deployment") -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SKY
    line.line.color.rgb = SKY
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(5.0), Inches(0.2))
    p = footer.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY


def add_bullets(slide, left: float, top: float, width: float, height: float, items: list[str], level: int = 0) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = level
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(5)
        p.bullet = True


def add_card(slide, left: float, top: float, width: float, height: float, title: str, body: list[str], fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = SKY
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    for item in body:
        child = tf.add_paragraph()
        child.text = item
        child.font.size = Pt(11)
        child.font.color.rgb = DARK
        child.space_before = Pt(4)


def add_table_slide(slide, left: float, top: float, width: float, height: float, headers: list[str], rows: list[list[str]]) -> None:
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for index, header in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_index % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK


def add_process_boxes(slide, steps: list[tuple[str, str]]) -> None:
    x = 0.7
    y = 1.8
    width = 3.9
    height = 1.1
    gap_x = 0.25
    gap_y = 0.25
    for index, (title, body) in enumerate(steps):
        row = index // 3
        col = index % 3
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x + col * (width + gap_x)),
            Inches(y + row * (height + gap_y)),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = SKY
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE
        child = tf.add_paragraph()
        child.text = body
        child.font.size = Pt(10)
        child.font.color.rgb = DARK


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "AgentCert Azure Production Deployment", "End-to-End Architecture and Component Stack", "1 / 17")
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.4))
    hero.fill.solid()
    hero.fill.fore_color.rgb = LIGHT_BLUE
    hero.line.color.rgb = SKY
    tf = hero.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "AgentCert"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "AI Agent Benchmarking for Chaos Engineering"
    p2.font.size = Pt(20)
    p2.font.color.rgb = DARK
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "April 2026"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Deployment Architecture Zones", None, "2 / 17")
    add_card(slide, 0.7, 1.6, 3.7, 1.3, "User Entry Zone", ["Azure DNS", "Application Gateway + WAF", "AKS ingress"], LIGHT_BLUE)
    add_card(slide, 4.8, 1.6, 3.7, 1.3, "Compute Zone", ["AKS cluster", "Control, execution, and benchmark workloads"], LIGHT_GREEN)
    add_card(slide, 8.9, 1.6, 3.7, 1.3, "Azure Services Zone", ["ACR, Key Vault, MongoDB/Cosmos DB", "Azure OpenAI, Storage, Monitor"], LIGHT_PURPLE)
    add_card(slide, 3.0, 3.3, 7.0, 1.25, "Observability Zone", ["Langfuse", "Azure Monitor", "OTEL-linked evidence path"], LIGHT_ORANGE)
    note = slide.shapes.add_textbox(Inches(1.1), Inches(5.2), Inches(11.0), Inches(1.0))
    p = note.text_frame.paragraphs[0]
    p.text = "Key point: Users enter through a secured edge, land on AKS, and rely on managed Azure services for secrets, models, telemetry, and state."
    p.font.size = Pt(18)
    p.font.color.rgb = DARK
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Detailed Azure Deployment Topology", "Customer-hosted platform inside Azure region with managed services on the right rail", "3 / 17")
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.35), Inches(2.25), Inches(4.95))
    left.fill.solid(); left.fill.fore_color.rgb = LIGHT_BLUE; left.line.color.rgb = SKY
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.95), Inches(1.05), Inches(7.2), Inches(5.25))
    center.fill.solid(); center.fill.fore_color.rgb = WHITE; center.line.color.rgb = SKY
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(1.35), Inches(2.4), Inches(4.95))
    right.fill.solid(); right.fill.fore_color.rgb = LIGHT_PURPLE; right.line.color.rgb = SKY
    add_card(slide, 0.62, 1.7, 1.9, 0.9, "Entry", ["Users", "DNS", "App Gateway", "Ingress"], LIGHT_BLUE)
    title = slide.shapes.add_textbox(Inches(3.2), Inches(1.3), Inches(3.0), Inches(0.3))
    p = title.text_frame.paragraphs[0]; p.text = "AKS Cluster"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BLUE
    add_card(slide, 3.2, 1.8, 2.1, 2.2, "litmus-chaos", ["Frontend", "Auth server", "GraphQL", "Registry APIs"], LIGHT_BLUE)
    add_card(slide, 5.55, 1.8, 2.1, 2.2, "litmus-exp", ["subscriber", "event-tracker", "workflow-controller", "chaos-operator"], LIGHT_GREEN)
    add_card(slide, 7.9, 1.8, 1.95, 2.2, "Agent namespaces", ["flash-agent", "install-agent", "target apps", "sidecar"], LIGHT_ORANGE)
    add_card(slide, 3.3, 4.35, 6.35, 0.9, "Shared integrations", ["Prometheus MCP", "Kubernetes MCP", "LiteLLM path", "RBAC and rollout policies"], LIGHT_GRAY)
    add_card(slide, 10.6, 1.6, 2.0, 3.9, "Azure managed services", ["ACR", "Key Vault", "MongoDB / Cosmos DB", "Azure OpenAI", "Langfuse", "Monitor + Storage"], LIGHT_PURPLE)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Edge Layer: DNS, WAF, Ingress", None, "4 / 17")
    add_card(slide, 0.8, 2.2, 2.8, 1.0, "Azure DNS", ["Production domain resolution"], LIGHT_BLUE)
    add_card(slide, 4.0, 2.2, 3.2, 1.0, "Application Gateway + WAF", ["TLS, DDoS and L7 routing"], LIGHT_PURPLE)
    add_card(slide, 7.7, 2.2, 2.8, 1.0, "AKS Ingress", ["Service routing into cluster"], LIGHT_GREEN)
    add_bullets(slide, 0.9, 4.0, 10.8, 2.0, [
        "Azure DNS exposes the production entrypoint.",
        "Application Gateway terminates TLS and enforces WAF policies.",
        "Ingress routes UI, API, and health traffic to the correct backend services.",
        "Recommended hardening: private connectivity to stateful backends and path-based routing for API isolation.",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Core Compute: AKS Cluster", None, "5 / 17")
    add_table_slide(slide, 0.55, 1.4, 12.1, 4.7,
        ["Pod / Service", "Namespace", "Technology", "Role"],
        [
            ["Frontend", "litmus-chaos", "React / TypeScript", "Web UI"],
            ["GraphQL API", "litmus-chaos", "Go / gqlgen / Gin", "Core API and orchestration"],
            ["Auth Service", "litmus-chaos", "Go / REST + gRPC", "JWT auth and user access"],
            ["Event Tracker", "litmus-exp", "Go controllers", "Infrastructure event streaming"],
            ["LitmusChaos + Argo", "litmus-exp", "Workflow engine", "Chaos orchestration"],
            ["Agents", "Agent namespaces", "Helm + agent framework", "Remediation execution"],
            ["Agent Sidecar", "Agent namespaces", "Python", "Trace injection and metadata"],
            ["LiteLLM Proxy", "Shared integration", "Python / FastAPI", "LLM gateway with callbacks"],
            ["OTEL Exporters", "Control-plane services", "Go / OTEL SDK", "Trace bridge to Langfuse"],
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Data, Secrets and Storage", None, "6 / 17")
    add_card(slide, 0.75, 1.8, 3.7, 2.1, "MongoDB / Cosmos DB for MongoDB", ["Users, projects, environments", "Agents, apps, experiments, runs", "Registry and audit state"], LIGHT_GREEN)
    add_card(slide, 4.85, 1.8, 3.3, 2.1, "Azure Key Vault", ["DB credentials", "JWT secrets", "Langfuse keys", "Azure OpenAI access"], LIGHT_BLUE)
    add_card(slide, 8.55, 1.8, 3.7, 2.1, "Azure Storage", ["Workflow artifacts", "Exports", "Backups and long-term evidence"], LIGHT_ORANGE)
    add_bullets(slide, 0.9, 4.5, 11.5, 1.4, [
        "Stateful services should use managed identity, secret rotation, and private endpoints where available.",
        "MongoDB keeps platform state while Langfuse stores deep LLM trace and evaluation payloads.",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "MongoDB Data Model", None, "7 / 17")
    for idx, (title, body, fill) in enumerate([
        ("project", ["_id", "name", "members[]", "state"], LIGHT_BLUE),
        ("environment", ["environment_id", "project_id", "type", "infra_ids[]"], LIGHT_GREEN),
        ("chaosInfrastructures", ["infra_id", "project_id", "name", "status"], LIGHT_ORANGE),
        ("apps_registrations", ["appId", "projectId", "environmentId", "namespace"], LIGHT_PURPLE),
        ("agentRegistry", ["agentId", "projectId", "vendor", "capabilities[]", "endpoint"], LIGHT_BLUE),
        ("chaosExperiments", ["experiment_id", "project_id", "infra_id", "revision[]"], LIGHT_GREEN),
        ("chaosExperimentRuns", ["experiment_run_id", "experiment_id", "phase", "resiliency_score"], LIGHT_ORANGE),
        ("chaosProbes", ["probe definitions", "fault mappings", "audit fields"], LIGHT_PURPLE),
        ("Supporting collections", ["chaosHubs", "imageRegistry", "serverConfig", "gitops", "faultStudios"], LIGHT_BLUE),
    ]):
        row = idx // 3
        col = idx % 3
        add_card(slide, 0.55 + col * 4.15, 1.4 + row * 1.55, 3.75, 1.3, title, body, fill)
    add_bullets(slide, 0.75, 6.3, 12.0, 0.6, [
        "Relationships: project -> environment -> apps_registrations; project -> agentRegistry; project -> chaosExperiments -> chaosExperimentRuns; chaosInfrastructures links experiments and runs.",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "LLM Integration and Observability", None, "8 / 17")
    add_card(slide, 1.0, 1.9, 2.2, 1.1, "Agent / LLM Call", ["Prompt, tool use, response"], WHITE)
    add_card(slide, 4.1, 1.9, 2.2, 1.1, "LiteLLM Proxy", ["Gateway, retries, callbacks"], LIGHT_GREEN)
    add_card(slide, 7.2, 1.9, 2.2, 1.1, "Azure OpenAI", ["Model inference"], LIGHT_ORANGE)
    add_card(slide, 10.0, 1.9, 2.2, 1.1, "Langfuse", ["Trace and score backend"], LIGHT_PURPLE)
    add_bullets(slide, 0.85, 3.7, 11.8, 2.2, [
        "LiteLLM centralizes model access, retries, auth abstraction, and callbacks.",
        "Langfuse stores traces, observations, evaluation outputs, and dashboards.",
        "OTEL bridge links Langfuse traces back to control-plane span context.",
        "Agent sidecar stabilizes trace IDs and carries agent metadata through benchmark runs.",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Detailed Deployment and Runtime Flow", None, "9 / 17")
    add_process_boxes(slide, [
        ("1. Build and publish", "Build GraphQL, auth, sidecar, LiteLLM, flash-agent, and install-agent images and push to ACR."),
        ("2. Bootstrap cluster", "Create namespaces, RBAC, CRDs, service accounts, and Litmus controllers before rollout changes."),
        ("3. Deploy control plane", "Roll out frontend, auth, and GraphQL first so registry and operator APIs are available."),
        ("4. Activate execution plane", "Bring up subscriber, event-tracker, workflow-controller, chaos-operator, and exporter services."),
        ("5. Run benchmark", "Link agent, target application, and fault plan; measure detection and remediation behavior."),
        ("6. Capture evidence", "Correlate OTEL, Langfuse, Kubernetes logs, and Mongo-backed run state into verdicts."),
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Container Registry: Azure ACR", None, "10 / 17")
    add_bullets(slide, 0.9, 1.6, 11.2, 4.7, [
        "agentcert/litmusportal-server:3.0.0",
        "agentcert/graphql-server:latest",
        "agentcert/auth-service:latest",
        "agentcert/agentcert-flash-agent:latest",
        "agentcert/agentcert-install-agent:latest",
        "agentcert/agent-sidecar:latest",
        "agentcert/agentcert-litellm-proxy:latest",
        "litmuschaos/* supporting images",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Azure Monitor and Logging", None, "11 / 17")
    add_card(slide, 0.9, 1.8, 5.5, 2.0, "Azure Monitor", ["Pod metrics", "Container logs", "Custom dashboards", "Alert rules"], LIGHT_BLUE)
    add_card(slide, 6.9, 1.8, 5.5, 2.0, "Langfuse", ["LLM traces", "Evaluation scores", "Agent decisions", "Performance trends"], LIGHT_PURPLE)
    add_bullets(slide, 0.9, 4.4, 11.4, 1.8, [
        "Monitor API latency, pod restarts, rollout health, queue depth, and trace quality.",
        "Operationally important signals include ReplicaSet creation, controller restart counts, and namespace quota pressure.",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Build and Release Pipeline", None, "12 / 17")
    add_process_boxes(slide, [
        ("Code push", "Source changes land in GitHub."),
        ("CI build", "GitHub Actions or Azure DevOps builds and tests images."),
        ("Push to ACR", "Versioned images are published to Azure Container Registry."),
        ("Deploy to AKS", "Helm and kubectl apply runtime configuration and chart updates."),
        ("Secret injection", "Runtime secrets come from Key Vault-backed flows."),
        ("Health validation", "Smoke tests, rollout checks, and post-deploy verification gates."),
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Security and Optional Hardening", None, "13 / 17")
    add_bullets(slide, 0.9, 1.7, 11.4, 4.8, [
        "Implemented: Key Vault for secrets management, ACR image scanning, AKS network policies.",
        "Optional hardening: VNet private endpoints for MongoDB, Key Vault, and Azure OpenAI.",
        "Optional hardening: NSGs, Microsoft Entra ID, pod security standards, and managed identity everywhere possible.",
        "Separate node pools for control-plane pods, experiment controllers, and heavier agent workloads when resource shapes diverge.",
    ])
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Operations, Scaling and Recovery", None, "14 / 17")
    add_card(slide, 0.8, 1.8, 3.8, 2.1, "Scaling model", ["Scale UI/API separately from execution controllers", "Use dedicated node pools for heavier benchmark workloads"], LIGHT_BLUE)
    add_card(slide, 4.8, 1.8, 3.8, 2.1, "Recovery model", ["Preserve MongoDB, Langfuse, Helm values, and exports", "Support replay and audit after outage"], LIGHT_GREEN)
    add_card(slide, 8.8, 1.8, 3.8, 2.1, "Support model", ["Runbooks for namespace bootstrap, rollout failures, ACR pulls, and secret access", "Day-2 checks for quotas and certificate expiry"], LIGHT_ORANGE)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Production Deployment Checklist", None, "15 / 17")
    add_card(slide, 0.75, 1.6, 3.8, 3.8, "Pre-deployment", ["Azure subscription and resource group", "AKS cluster", "MongoDB / Cosmos DB", "ACR", "Key Vault", "Azure OpenAI"], LIGHT_BLUE)
    add_card(slide, 4.8, 1.6, 3.2, 3.8, "Deployment", ["Build and push images", "Configure Helm values", "Sync secrets", "Run Helm install or upgrade", "Pass health checks"], LIGHT_GREEN)
    add_card(slide, 8.35, 1.6, 4.1, 3.8, "Post-deployment", ["DNS points to App Gateway", "TLS certificate active", "Monitoring dashboards live", "Backups enabled", "DR plan documented"], LIGHT_ORANGE)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Production Resource Estimates", None, "16 / 17")
    add_table_slide(slide, 0.7, 1.5, 11.9, 4.8,
        ["Component", "SKU / Configuration", "Monthly Cost"],
        [
            ["AKS Cluster", "3 nodes, Standard_D2s_v3", "$300-400"],
            ["Cosmos DB MongoDB", "400 RU/s, multi-region option", "$200-500"],
            ["Azure OpenAI", "gpt-4.1-mini + embeddings", "$50-500 usage-based"],
            ["Azure Container Registry", "Premium tier", "$100-200"],
            ["Key Vault", "Standard tier", "$10"],
            ["Azure Monitor", "Ingestion + 30-day retention", "$50-200"],
            ["App Gateway", "Standard tier with WAF", "$30-50"],
            ["Storage Account", "Blob, geo-redundant", "$20-100"],
            ["Total baseline", "Varies by region and load", "~$750-2000/month"],
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Summary and Next Steps", None, "17 / 17")
    add_bullets(slide, 0.8, 1.6, 11.6, 2.8, [
        "Detailed Azure deployment topology with namespace separation.",
        "MongoDB data model for platform state and benchmark evidence.",
        "Secure edge, AKS workloads, Azure managed services, observability, and CI/CD coverage.",
        "Operational scaling, recovery, and day-2 support guidance.",
    ])
    add_bullets(slide, 0.8, 4.5, 11.6, 1.7, [
        "Provision Azure resources.",
        "Configure production Helm values.",
        "Enable dashboards and alert rules.",
        "Run staging dry-run deployment and rollback rehearsal.",
        "Train operations team on runbooks.",
    ])
    add_footer(slide)

    return prs


def main() -> None:
    prs = build_deck()
    prs.save(str(OUTPUT_PATH))
    print(f"Created {OUTPUT_PATH}")


if __name__ == "__main__":
    main()